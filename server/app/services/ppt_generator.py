import json
import os
from typing import Dict, Any
from pptx import Presentation
from openai import OpenAI

from app.config import DEEPSEEK_API_KEY, SERVER_DIR
from app.core.logger import setup_logger
from app.schemas.ppt_generator import PPTGenerationRequest, PPTGenerationResponse, PPTSlide

# 设置日志
logger = setup_logger("ppt_generator_service")

# PPT文件存储目录
PPT_FILES_DIR = SERVER_DIR / "app" / "documents" / "ppt_files"
PPT_FILES_DIR.mkdir(exist_ok=True, parents=True)

# 创建 API 客户端
client = OpenAI(
    api_key=DEEPSEEK_API_KEY,
    base_url="https://api.deepseek.com"
)



async def generate_ppt_content(request: PPTGenerationRequest) -> Dict[str, Any]:
    """
    调用Deepseek API生成PPT内容
    """
    prompt = f"""你是一位经验丰富的教师和课件专家，请为以下教学内容设计一个详细的PPT:
    
    标题: {request.title}
    学科: {request.subject}
    目标年级: {request.target_grade}
    教学目标: {request.teaching_target}
    教学重点: {', '.join(request.key_points)}
    {f"其他信息: {request.additional_info}" if request.additional_info else ""}
    
    请生成{request.slide_count}张幻灯片内容，包括:
    1. 封面页(标题+简介)
    2. 学习目标页
    3. 内容页(核心知识讲解)
    4. 案例/示例页
    5. 练习/活动页
    6. 总结页
    
    请以JSON格式返回，格式如下:
    {{
        "title": "PPT标题",
        "slides": [
            {{
                "title": "幻灯片标题",
                "content": "幻灯片内容(支持Markdown格式)",
                "note": "教师备注说明"
            }}
        ]
    }}
    
    只返回JSON格式结果，不要有其他说明文字。
    """

    try:
        # 使用 OpenAI 客户端调用 API
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": "你是一个专业的教育资源制作助手，擅长生成教学PPT内容。"},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=4000
        )

        content = response.choices[0].message.content

        # 从返回的内容中提取JSON
        try:
            # 确保返回的是纯JSON
            if content.startswith("```json"):
                content = content.split("```json")[1].split("```")[0].strip()
            elif content.startswith("```"):
                content = content.split("```")[1].split("```")[0].strip()

            ppt_data = json.loads(content)
            return ppt_data
        except json.JSONDecodeError as e:
            logger.error(f"解析JSON失败: {e}\n内容: {content}")
            raise Exception("解析生成的PPT内容失败")

    except Exception as e:
        logger.error(f"生成PPT内容时出错: {str(e)}")
        raise Exception(f"生成PPT内容失败: {str(e)}")


async def create_pptx_file(ppt_data: Dict[str, Any], file_name: str) -> str:
    """
    根据生成的内容创建PPTX文件

    Args:
        ppt_data: PPT内容数据
        file_name: 完整的文件名(包含扩展名)

    Returns:
        生成文件的完整路径
    """
    try:
        # 创建演示文稿
        prs = Presentation()

        # 添加幻灯片
        for slide_data in ppt_data["slides"]:
            # 使用标题和内容布局
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            # 设置标题
            title = slide.shapes.title
            title.text = slide_data["title"]

            # 设置内容
            content_shape = slide.placeholders[1]
            content_shape.text = slide_data["content"]

            # 如果有备注，添加备注
            if "note" in slide_data and slide_data["note"]:
                slide.notes_slide.notes_text_frame.text = slide_data["note"]

        # 保存文件
        file_path = PPT_FILES_DIR / file_name
        prs.save(str(file_path))

        return str(file_path)

    except Exception as e:
        logger.error(f"创建PPTX文件时出错: {str(e)}")
        raise Exception(f"创建PPTX文件失败: {str(e)}")


async def generate_ppt(request: PPTGenerationRequest, staff_id: str) -> PPTGenerationResponse:
    """
    生成PPT的主要服务函数
    """
    # 生成唯一标识
    request_id = os.urandom(4).hex()

    # 调用API生成PPT内容
    ppt_data = await generate_ppt_content(request)

    # 创建幻灯片列表
    slides = [PPTSlide(
        title=slide["title"],
        content=slide["content"],
        note=slide.get("note")
    ) for slide in ppt_data["slides"]]

    # 生成文件名，包含用户标识
    file_name = f"user_{staff_id}_{request_id}_{ppt_data['title'].replace(' ', '_')}.pptx"

    # 创建PPTX文件
    file_path = await create_pptx_file(ppt_data, file_name)

    # 返回结果
    return PPTGenerationResponse(
        title=ppt_data["title"],
        slides=slides,
        file_path=file_path
    )