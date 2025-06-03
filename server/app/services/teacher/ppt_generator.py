import uuid
from datetime import datetime
from typing import Dict, Any, List

from openai import OpenAI
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.config import DEEPSEEK_API_KEY, SERVER_DIR
from app.core.logger import setup_logger
from app.schemas.ppt_generator import (
    PPTGenerationRequest, PPTGenerationResponse, PPTSlide,
    PPTOutlineResponse, PPTGenerationFromOutlineRequest
)

# 设置日志
logger = setup_logger("ppt_generator_service")

# PPT文件存储目录
PPT_FILES_DIR = SERVER_DIR / "app" / "documents" / "ppt_files"
PPT_OUTLINE_DIR = SERVER_DIR / "app" / "documents" / "ppt_outlines"
PPT_FILES_DIR.mkdir(exist_ok=True, parents=True)
PPT_OUTLINE_DIR.mkdir(exist_ok=True, parents=True)

# 创建 API 客户端
client = OpenAI(
    api_key=DEEPSEEK_API_KEY,
    base_url="https://api.deepseek.com"
)


async def generate_ppt_outline(request: PPTGenerationRequest, staff_id: str) -> PPTOutlineResponse:
    """
    生成PPT的Markdown格式大纲
    """
    # 生成唯一标识
    request_id = uuid.uuid4().hex[:8]

    logger.info(f"开始生成PPT大纲: 标题={request.title}, ID={request_id}")

    prompt = f"""你是一位经验丰富的教师和课件专家，请为以下教学内容设计一个详细且内容丰富的PPT大纲:
    
    标题: {request.title}
    学科: {request.subject}
    目标年级: {request.target_grade}
    教学目标: {request.teaching_target}
    教学重点: {', '.join(request.key_points)}
    幻灯片数量: {request.slide_count}
    {f"其他信息: {request.additional_info}" if request.additional_info else ""}
    
    请生成{request.slide_count}张幻灯片的md格式设计，要求每张幻灯片的内容具体、丰富、实用，避免空洞的表述。内容需要包括:
    1. 封面页(标题+简介)
    2. 学习目标页(详细列出具体可衡量的学习成果)
    3. 内容页(核心知识讲解，包含具体概念、定义、原理和示例代码)
    4. 案例/示例页(真实可执行的代码示例，而非空泛描述)
    5. 练习/活动页(有明确步骤和要求的练习题)
    6. 总结页(具体的知识点总结，不要简单重复前面内容)
    
    对于编程相关内容，请提供实际的代码片段而非仅描述代码功能。
    对于概念解释，请给出明确的定义和具体的例子。
    对于操作步骤，请提供详细的分步骤说明。
    
    请提供一个完整的Markdown格式大纲，包含每张幻灯片的标题和内容。
    
    请确保生成了符合要求的幻灯片数量，字数适中且内容有教学价值。
    """

    try:
        # 使用 OpenAI 客户端调用 API
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": "你是一个专业的教育资源制作助手，擅长生成教学PPT内容。"},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=8000,
            stream=False
        )

        md_content = response.choices[0].message.content
        logger.info(f"成功从API获取大纲内容，请求ID={request_id}")

        # 保存大纲到文件
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        outline_path = PPT_OUTLINE_DIR / f"outline_{staff_id}_{timestamp}_{request.title}.md"
        with open(outline_path, "w", encoding="utf-8") as f:
            f.write(md_content)

        logger.info(f"大纲已保存到文件: {outline_path}")

        return PPTOutlineResponse(
            request_id=request_id,
            title=request.title,
            outline_md=md_content,
        )

    except Exception as e:
        logger.error(f"生成PPT大纲失败: {str(e)}")
        raise Exception(f"生成PPT大纲失败: {str(e)}")


async def generate_ppt_from_outline(
        request: PPTGenerationFromOutlineRequest,
        staff_id: str
) -> PPTGenerationResponse:
    """
    第二步：从修改后的大纲生成PPT
    """
    logger.info(f"从大纲生成PPT: 请求ID={request.request_id}, 教师ID={staff_id}")

    # 检查大纲文件是否存在
    outline_path = PPT_OUTLINE_DIR / f"outline_{staff_id}_{request.request_id}.md"
    if not outline_path.exists():
        logger.error(f"大纲文件不存在: {outline_path}")
        raise Exception("大纲文件不存在")

    try:
        # 将大纲解析为结构化数据
        structured_data = parse_outline_to_json(request.outline_md, request.title)

        # 生成文件名，包含教师工号以实现权限隔离
        file_name = f"teacher_{staff_id}_{request.request_id}_{request.title.replace(' ', '_')}.pptx"

        # 创建高质量PPT
        file_path = await create_enhanced_pptx(structured_data, file_name, request.design_preference)

        # 构建返回结果
        slides = [
            PPTSlide(title=slide["title"], content=slide["content"], note=slide.get("note"))
            for slide in structured_data["slides"]
        ]

        return PPTGenerationResponse(
            title=structured_data["title"],
            slides=slides,
            file_path=file_path
        )

    except Exception as e:
        logger.error(f"从大纲生成PPT失败: {str(e)}")
        raise Exception(f"从大纲生成PPT失败: {str(e)}")


def parse_outline_to_json(outline_md: str, title: str) -> Dict[str, Any]:
    """解析Markdown大纲为JSON结构"""
    result = {"title": title, "slides": []}

    try:
        lines = outline_md.split("\n")
        current_slide = None
        notes = {}
        in_notes_section = False

        for line in lines:
            line = line.strip()

            # 处理标题行
            if line.startswith("## 幻灯片"):
                # 如果有当前处理的幻灯片，保存它
                if current_slide:
                    result["slides"].append(current_slide)

                # 提取幻灯片标题
                slide_title = line.split("：", 1)[1] if "：" in line else line[2:].strip()
                current_slide = {"title": slide_title, "content": "", "note": ""}

            # 处理教师备注部分
            elif line.startswith("## 教师备注"):
                in_notes_section = True
                # 保存最后一个幻灯片
                if current_slide:
                    result["slides"].append(current_slide)
                    current_slide = None

            # 处理备注内容
            elif in_notes_section and line.startswith("- 幻灯片"):
                parts = line.split("：", 1)
                if len(parts) == 2:
                    slide_num = parts[0].replace("- 幻灯片", "").strip()
                    try:
                        slide_idx = int(slide_num) - 1
                        if 0 <= slide_idx < len(result["slides"]):
                            result["slides"][slide_idx]["note"] = parts[1].strip()
                    except ValueError:
                        pass

            # 处理幻灯片内容
            elif current_slide and line:
                if line.startswith("- 标题："):
                    current_slide["title"] = line[5:].strip()
                elif line.startswith("- 内容："):
                    current_slide["content"] = ""
                else:
                    current_slide["content"] += line + "\n"

        # 确保最后一个幻灯片被添加
        if current_slide and not in_notes_section:
            result["slides"].append(current_slide)

        return result

    except Exception as e:
        logger.error(f"解析大纲失败: {str(e)}")
        raise Exception(f"解析大纲失败: {str(e)}")


async def create_enhanced_pptx(data: Dict[str, Any], file_name: str, design_preference: str = None) -> str:
    """创建增强版PPT，更美观丰富"""
    try:
        # 创建演示文稿
        prs = Presentation()

        # 设置主题颜色
        colors = get_theme_colors(design_preference)

        for i, slide_data in enumerate(data["slides"]):
            # 根据内容选择合适的布局
            if i == 0:  # 封面页使用标题页布局
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)

                # 设置标题
                title = slide.shapes.title
                title.text = slide_data["title"]
                format_text_frame(title.text_frame, font_size=44, bold=True, color=colors["title"])

                # 设置副标题
                subtitle = slide.placeholders[1]
                subtitle.text = extract_subtitle(slide_data["content"])
                format_text_frame(subtitle.text_frame, font_size=28, color=colors["subtitle"])

            else:
                # 根据内容选择布局
                content = slide_data["content"]
                if "* " in content and content.count("* ") > 3:
                    # 使用带项目符号的布局
                    slide_layout = prs.slide_layouts[1]  # Title and Content
                elif i == len(data["slides"]) - 1:
                    # 最后一页用总结布局
                    slide_layout = prs.slide_layouts[5]  # Title Only
                else:
                    # 普通内容页
                    slide_layout = prs.slide_layouts[2]  # Section Header

                slide = prs.slides.add_slide(slide_layout)

                # 设置标题
                title = slide.shapes.title
                title.text = slide_data["title"]
                format_text_frame(title.text_frame, font_size=36, bold=True, color=colors["title"])

                # 设置内容
                if len(slide.placeholders) > 1:
                    content_shape = slide.placeholders[1]
                    format_content(content_shape, slide_data["content"], colors)
                else:
                    # 创建自定义文本框
                    left = Inches(1)
                    top = Inches(2)
                    width = Inches(8)
                    height = Inches(4)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.text = slide_data["content"]
                    format_text_frame(tf, font_size=24, color=colors["text"])

            # 添加页脚
            add_footer(slide, i + 1, len(data["slides"]), data["title"])

            # 添加备注
            if slide_data.get("note"):
                slide.notes_slide.notes_text_frame.text = slide_data["note"]

        # 保存文件
        file_path = PPT_FILES_DIR / file_name
        prs.save(str(file_path))

        logger.info(f"PPT文件已保存: {file_path}")
        return str(file_path)

    except Exception as e:
        logger.error(f"创建PPT失败: {str(e)}")
        raise Exception(f"创建PPT失败: {str(e)}")


def format_text_frame(text_frame, font_size=24, bold=False, italic=False, color=None, alignment=PP_ALIGN.LEFT):
    """格式化文本框架"""
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.alignment = alignment

    run = p.add_run()
    run.text = text_frame.text

    font = run.font
    font.size = Pt(font_size)
    font.bold = bold
    font.italic = italic

    if color:
        font.color.rgb = color


def format_content(shape, content: str, colors: Dict):
    """格式化内容，支持Markdown格式"""
    tf = shape.text_frame
    tf.clear()

    lines = content.split("\n")
    first = True

    for line in lines:
        line = line.strip()
        if not line:
            continue

        # 添加段落
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False

        # 设置段落格式
        if line.startswith("* ") or line.startswith("- "):
            p.level = 1
            line = line[2:]
        elif line.startswith("  * ") or line.startswith("  - "):
            p.level = 2
            line = line[4:]

        # 添加文本
        run = p.add_run()
        run.text = line

        # 设置字体
        font = run.font
        font.size = Pt(24 if p.level == 0 else 20 if p.level == 1 else 18)

        # 设置颜色
        if p.level == 0:
            font.color.rgb = colors["text"]
        else:
            font.color.rgb = colors["bullet"]
            font.bold = p.level == 1  # 第一级项目符号加粗


def add_footer(slide, slide_number, total_slides, title):
    """添加页脚"""
    left = Inches(0.3)
    top = Inches(6.8)
    width = Inches(9.4)
    height = Inches(0.3)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = f"{title} | 第 {slide_number}/{total_slides} 页"

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT

    run = p.runs[0]
    font = run.font
    font.size = Pt(12)
    font.color.rgb = RGBColor(128, 128, 128)


def extract_subtitle(content: str) -> str:
    """从内容中提取适合作为副标题的文本"""
    lines = content.split("\n")
    clean_lines = [line.strip() for line in lines if line.strip() and not line.strip().startswith("*")]

    return " | ".join(clean_lines[:2]) if clean_lines else ""


def get_theme_colors(preference: str = None) -> Dict[str, RGBColor]:
    """根据设计偏好获取主题颜色"""
    themes = {
        "default": {
            "title": RGBColor(31, 73, 125),  # 深蓝色
            "subtitle": RGBColor(68, 114, 196),  # 蓝色
            "text": RGBColor(0, 0, 0),  # 黑色
            "bullet": RGBColor(89, 89, 89)  # 深灰色
        },
        "modern": {
            "title": RGBColor(0, 102, 153),  # 蓝绿色
            "subtitle": RGBColor(0, 153, 204),  # 亮蓝色
            "text": RGBColor(51, 51, 51),  # 深灰色
            "bullet": RGBColor(0, 153, 153)  # 青色
        },
        "warm": {
            "title": RGBColor(204, 51, 0),  # 红色
            "subtitle": RGBColor(231, 145, 36),  # 橙色
            "text": RGBColor(51, 51, 0),  # 棕色
            "bullet": RGBColor(153, 51, 0)  # 深红色
        },
        "elegant": {
            "title": RGBColor(68, 35, 102),  # 紫色
            "subtitle": RGBColor(109, 71, 150),  # 淡紫色
            "text": RGBColor(0, 0, 0),  # 黑色
            "bullet": RGBColor(76, 47, 107)  # 深紫色
        }
    }

    if preference and preference in themes:
        return themes[preference]
    else:
        return themes["default"]
